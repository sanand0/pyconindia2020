from pptx import Presentation

prs = Presentation('input.pptx')
prs.slides[0].shapes[0].text = 'Hello world'
prs.save('output.pptx')
